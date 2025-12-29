from io import BytesIO
from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image
import csv


def parse_file_from_bytes(file_bytes: bytes, file_name: str) -> str:
    """
    Universal file parser that extracts text from PDF, DOCX, XLSX, CSV, PPTX, TXT.
    OCR is disabled in serverless environments.
    """
    ext = file_name.lower().split(".")[-1]

    # PDF
    if ext == "pdf":
        reader = PdfReader(BytesIO(file_bytes))
        return "".join(page.extract_text() or "" for page in reader.pages)

    # DOCX
    elif ext == "docx":
        doc = Document(BytesIO(file_bytes))
        return "\n".join(p.text for p in doc.paragraphs)

    # TXT
    elif ext == "txt":
        return file_bytes.decode("utf-8", errors="ignore")

    # XLSX
    elif ext == "xlsx":
        wb = load_workbook(BytesIO(file_bytes), data_only=True)
        text = ""
        for sheet in wb:
            for row in sheet.iter_rows(values_only=True):
                text += " ".join(str(cell) for cell in row if cell is not None) + "\n"
        return text

    # CSV
    elif ext == "csv":
        f = BytesIO(file_bytes).read().decode("utf-8", errors="ignore").splitlines()
        return "\n".join(" ".join(row) for row in csv.reader(f))

    # PPTX
    elif ext == "pptx":
        prs = Presentation(BytesIO(file_bytes))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    # Images (OCR disabled)
    elif ext in ["png", "jpg", "jpeg"]:
        return "[Image OCR is not supported in this environment]"

    else:
        return f"[Unsupported file format: {ext}]"
